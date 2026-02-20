"""
PowerPoint Presentation Loader.

Extracts text from .pptx files including:
- Slide title and body text
- Speaker notes (often contain more detail than the slides themselves)
- Table content from slides

Strategy: Each slide becomes a Document. We combine the slide text and
speaker notes because together they tell the full story.

Analogy: A slide deck is like a movie trailer (the slides) plus the full
script (speaker notes). For RAG, you want BOTH — the trailer gives
keywords, and the script gives context.
"""
from langchain_core.documents import Document
from pptx import Presentation
import logging

logger = logging.getLogger(__name__)


def load_pptx(file_path: str) -> list[Document]:
    """
    Load a PowerPoint presentation and return LangChain Document objects.

    Each slide becomes a separate Document containing:
    - Slide title
    - All text from shapes (bullets, text boxes)
    - Speaker notes
    - Table content

    Args:
        file_path: Path to the .pptx file

    Returns:
        List of Documents, one per slide
    """
    documents = []

    try:
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, start=1):
            parts = [f"Slide {slide_num}"]

            # Extract title
            if slide.shapes.title and slide.shapes.title.text.strip():
                parts.append(f"Title: {slide.shapes.title.text.strip()}")

            # Extract all text from shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text and text != (slide.shapes.title.text.strip() if slide.shapes.title else ""):
                            parts.append(text)

                # Extract tables on slides
                if shape.has_table:
                    table = shape.table
                    headers = [cell.text.strip() for cell in table.rows[0].cells]
                    for row in table.rows[1:]:
                        row_text = " | ".join(
                            f"{h}: {cell.text.strip()}"
                            for h, cell in zip(headers, row.cells)
                        )
                        parts.append(row_text)

            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Speaker Notes] {notes}")

            full_text = "\n".join(parts)

            if len(full_text.strip()) > 20:  # Skip near-empty slides
                documents.append(Document(
                    page_content=full_text.strip(),
                    metadata={
                        "source": file_path,
                        "file_type": "pptx",
                        "slide_number": slide_num,
                    }
                ))

        logger.info(f"Loaded PPTX: {file_path} — {len(documents)} slides")

    except Exception as e:
        logger.error(f"Error loading PPTX {file_path}: {e}")
        raise

    return documents
